#Python3

# Updates charts and maps in LAO Market Overview PowerPoints

import lao
import os
from pprint import pprint
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import sys
import fun_text_date as td
import win32com.client

# Select Market Menu
def selectMarketMenu() :
	lMarkets = [
		'ATL - Atlanta',
		'AUS - Austin',
		'BOI - Boise',
		'CLT - Charlotte',
		'DFW - Dallas Ft Worth',
		'GSP - Greenville',
		'HNT - Huntsville',
		'HOU - Houston',
		'JAX - Jacksonville',
		'KCI - Kansas City',
		'LVS - Las Vegas',
		'NAZ - Northern Arizona',
		'NSH - Nashville',
		'ORL - Orlando',
		'RNO - Reno',
		'SLC - Salt Lake City',
		'TPA - Tampa Bay/Sarasota',
		'TUC - Tucson']
	while 1:
		lao.banner('MO PowerPoint Chart Placement v04')
		index = 1
		for market in lMarkets:
			print(' {0:2}) {1}\n'.format(index, market))
			index += 1
		print('\n 00) Quit')
		ui = td.uInput('\n Select Market > ')
		if ui == '00':
			exit('\n Program terminated...')
		uMarket = lMarkets[int(ui) - 1]
		ui = td.uInput('\n Process {0} [0/1/00] > '.format(uMarket))
		if ui == '1':
			lMkt = uMarket.split(' - ')
			return lMkt[0]
		elif ui == '00':
			exit('\n Program terminated...')
		
# Set the file name of the PPTX to use
def setPPTXfilename():
		mkt = d['Market']
		print('\n Market: {0}'.format(d['Market']))
		pptxFile = 'LAO_{0}_Market_Overview_{1}.pptx'.format(d['Market'], curQtr)
		fin = '{0}{1}'.format(pptxPath, pptxFile)
		prs = Presentation(fin)
		return mkt, fin, prs

# Find the slide to update the picture baed on slide title
def findTextInSlide():
	td.colorText('\n Slide Title: {0}'.format(d['SlideTitle']), 'CYAN')
	for slide in prs.slides:
		foundSlide = False
		# Find the slide to update based on the text of the slide's title
		for shape in slide.shapes:
			if shape.has_text_frame:
			# Determine if text is in text frame
				if(shape.text.find(d['SlideTitle']))!=-1:
					slideIndex = prs.slides.index(slide)+1
					print(' Slide: {0}'.format(slideIndex))
					print(' {0} TRUE'.format(d['Chart']))
					return True, slide
	return False, None

# Cycle through the shapes to find the PICTURE type and replace it.
def replacePicture(slide):
	picture_file_name = os.path.basename(chart)
	print('\n Replace Picture: {0}'.format(picture_file_name))
	for shape in slide.shapes:
		if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
			pic = shape
			pic = pic._element
			pic.getparent().remove(pic) # delete the shape with the picture
			break
				
	# Add new chart or map picture to slide
	pic = slide.shapes.add_picture(chart, Inches(lf), Inches(tp), Inches(wd), Inches(ht))
	# Move picture to the back
	if d['SendToBack'] == 'Yes':
		slide.shapes._spTree.insert(2, pic._element)

# Replace subdivision legend counts text in subdivision slides
def replaceSubdivisionLegendCounts(slide):
	lLast = d['LastSubCount'].split(':')
	# print(lLast)
	lCur = d['CurSubCount'].split(':')
	print
	# print('here2')
	# print(lCur)
	# print(d['CurSubCount'])
	# loopSlideShapes(slide, strSearch, strReplace)
	index = 0
	# print('\n Dict ------------------------------')
	# pprint(d)
	# print('------------------------------------')
	for strSearch in lLast:
		
		strReplace = lCur[index]
		
		print(f'\n Index:           {index}')
		print(f' lLast strSearch: {strSearch}')
		print(f' lCur strReplace: {strReplace}')

		loopSlideShapes(slide, strSearch, strReplace)
		
		index += 1

# Loop through slide's shapes and grouped shapes
def loopSlideShapes(slide, strSearch, strReplace):
	# search and replace text in PowerPoint while preserving formatting
	group_shapes = [
			shp for shp in slide.shapes
			if shp.shape_type == MSO_SHAPE_TYPE.GROUP
		]
	for group_shape in group_shapes:
		for shape in group_shape.shapes:
			# print(' Group Shape replace text')
			replaceTextInShape(shape, strSearch, strReplace)
	# Loop through shapes
	for shape in slide.shapes:
		# print(' Shape replace text')
		replaceTextInShape(shape, strSearch, strReplace)

# Loop through shapes
def replaceTextInShape(shape, strSearch, strReplace):
	if shape.has_text_frame:
		# Determine if text to be replaced is in text frame
		if(shape.text.find(strSearch))!=-1:
			text_frame = shape.text_frame
			# Get number of lines (carrage returns) in text frame
			lenPar = len(text_frame.paragraphs) - 0
			i = 0
			# Loop through lines and replace text
			textFoundReplaced = False
			while 1:
				try:
					cur_text = text_frame.paragraphs[i].runs[0].text
				except IndexError:
					print('Failed {0}'.format(strSearch))
					i += 1
					if i >= lenPar:
						print(' break')
						break
					continue
				new_text = cur_text.replace(str(strSearch), str(strReplace))
				
				# Check if string already replaced
				if new_text == cur_text:
					# print(' EQUAL - cur_text: {0}   new_text: {1}'.format(cur_text, new_text))
					i += 1
					# Exit while loop
					if i >= lenPar:
						print(' break')
						break
					continue

				text_frame.paragraphs[i].runs[0].text = new_text
				textFoundReplaced = True
				td.colorText(f' Find Replace Success - cur_text: {cur_text}   new_text: {new_text}', 'GREEN')
				break
				# Exit while loop
				if i >= lenPar:
					print(' break')
					break
				else:
					i += 1
			if textFoundReplaced is False:
				td.warningMsg('Find Replace FAIL - cur_text: {0}   new_text: {1}'.format(cur_text, new_text))
				lao.holdup()

# Set Variables
curQtr = lao.getDateQuarter(lastquarter=True)
pptxPath = 'F:/Research Department/MIMO/PowerPoints/{0}/'.format(curQtr)
chartsPath = 'F:/Research Department/MIMO/PowerPoints/Template Data/'

while 1:
	dCharts = lao.spreadsheetToDict('F:/Research Department/MIMO/Spreadsheets/{0}/MO Python Chart Placement Database.xlsx'.format(curQtr))
	ui = selectMarketMenu()
	lao.banner('MO PowerPoint Chart Placement v04')
	mkt = 'None'
	for row in dCharts:
		foundSlide = False
		d = dCharts[row]

		# Skip all but selected market
		if d['Market'] != ui:
			continue
		
		# Skip Population chart since only updated once per year
		if d['Update'] == 'NO':
			continue

		# set pptx file name if switching to next market
		if mkt != d['Market']:
			mkt, fin, prs = setPPTXfilename()
			fout = fin.replace('.pptx', '_v1.pptx')

		# Set image parms: Left, Top, Width, Height parms
		lf, tp, wd, ht = float(d['Left']), float(d['Top']), float(d['Width']), float(d['Height'])

		# Charts
		if d['FileFormat'] == 'png':
			chart = '{0}{1} {2}.{3}'.format(chartsPath, d['Market'], d['Chart'], d['FileFormat'])
		# Arcmap maps
		else:
			chart = '{0}{1}.{2}'.format(chartsPath, d['Chart'], d['FileFormat'])

		# Cycle through slides
		foundSlide, slide = findTextInSlide()
		# If slide to updated then replace the chart or map
		if foundSlide:
			replacePicture(slide)
			# Replace subdivision legend counts text in subdivision slides
			if d['LastSubCount'] != None:
				print(' Replacing legend text')
				if mkt != 'TUC':
					replaceSubdivisionLegendCounts(slide)
			else:
				print(' No legend')

	print('\n Saving...')
	prs.save(fout) # save
	lao.openFile(fout)
	lao.holdup()







