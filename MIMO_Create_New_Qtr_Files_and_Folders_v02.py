# Copies inDesign files and creates new quarter's directory and files.

import lao
from glob import glob
import fun_text_date as td
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pprint import pprint
import shutil

# get Search and Replace List of lists
def get_Qtrs_lSandRs():
	curQtr = lao.getDateQuarter(lastquarter=True)
	intCurYear = int(curQtr[:4])
	intLastYear = intCurYear - 1
	intCurQtr = int(curQtr[5:])
	intLastQtr = intCurQtr - 1

	isYearChange = False	
	if intLastQtr == 0:
		intLastQtr = 4
		lastQtr = '{0}Q{1}'.format(intLastYear, intLastQtr)
	else:
		lastQtr = '{0}Q{1}'.format(intCurYear, intLastQtr)
	# Make search & replace quarter variables
	if intLastQtr == 1:
		srLastOrdinal, srCurOrdinal = '1st', '2nd'
		srLastMoYr = 'March 20'
		srCurMoYr = 'June 20'
		srLastMoYrShort = 'Mar 20'.format(intCurYear)
		srCurMoYrShort = 'Jun 20'
		srLastUnempTitle = 'March Unemployment Rate'
		srCurUnempTitle = 'June Unemployment Rate'
	elif intLastQtr == 2:
		# srLastOrdinal, srCurOrdinal = '2nd', '3rd'
		# srLastMoYr = 'June 20'
		# srCurMoYr = 'Sept. 20'
		# srLastMoYrShort = 'Jun 20'
		# srCurMoYrShort = 'Sep 20'
		# srLastUnempTitle = 'June Unemployment Rate'
		# srCurUnempTitle = 'September Unemployment Rate'
		srLastOrdinal, srCurOrdinal = '4th', '2nd'
		srLastMoYr = 'Dec. 2022'
		srCurMoYr = 'June 2023'
		srLastMoYrShort = 'Dec 2022'
		srCurMoYrShort = 'June 2023'
		srLastUnempTitle = 'December Unemployment Rate'
		srCurUnempTitle = 'June Unemployment Rate'
		# srLastEmpRateYears = '	2022	2023	Change'
		# srCurEmpRateYears = '	2022	2023	Change'
		srLastYear = ' 2022'
		srCurYear = ' 2023'
		isYearChange = True
	elif intLastQtr == 3:
		srLastOrdinal, srCurOrdinal = '3rd', '4th'
		srLastMoYr = 'Sept. 20'
		srCurMoYr = 'Dec. 20'
		srLastMoYrShort = 'Sep 20'
		srCurMoYrShort = 'Dec 20'
		srLastUnempTitle = 'September Unemployment Rate'
		srCurUnempTitle = 'December Unemployment Rate'
	elif intLastQtr == 4:
		srLastOrdinal, srCurOrdinal = '4th', '2nd'
		srLastMoYr = 'Dec. 2022'
		srCurMoYr = 'June 2023'
		srLastMoYrShort = 'Dec 2022'
		srCurMoYrShort = 'June 2023'
		srLastUnempTitle = 'December Unemployment Rate'
		srCurUnempTitle = 'June Unemployment Rate'
		# srLastEmpRateYears = '	2022	2023	Change'
		# srCurEmpRateYears = '	2022	2023	Change'
		srLastYear = ' 2022'
		srCurYear = ' 2023'
		isYearChange = True
	# Miss doing 1st qtr of 2023
	# elif intLastQtr == 4:
	# 	srLastOrdinal, srCurOrdinal = '4th', '1st'
	# 	srLastMoYr = 'Dec. 2021'
	# 	srCurMoYr = 'March 2022'
	# 	srLastMoYrShort = 'Dec 2021'
	# 	srCurMoYrShort = 'Mar 2022'
	# 	srLastUnempTitle = 'December Unemployment Rate'
	# 	srCurUnempTitle = 'March Unemployment Rate'
	# 	srLastEmpRateYears = '	2020	2021	Change'
	# 	srCurEmpRateYears = '	2021	2022	Change'
	# 	srLastYear = ' 2021'
	# 	srCurYear = ' 2022'
	# 	isYearChange = True

	lSandRs = [
		[srLastOrdinal, srCurOrdinal],
		[srLastMoYr, srCurMoYr],
		[srLastMoYrShort, srCurMoYrShort],
		[srLastUnempTitle, srCurUnempTitle]
		]
	
	# srlastUSUnemp = 'National:	6.7%	3.9%	2.8%'
	# srCurUSUnemp = 'National:	6.0%	3.6%	2.4%'
	# lSandRs = [
	# 	[srLastOrdinal, srCurOrdinal],
	# 	[srLastMoYr, srCurMoYr],
	# 	[srLastMoYrShort, srCurMoYrShort],
	# 	[srlastUSUnemp, srCurUSUnemp],
	# 	[srLastUnempTitle, srCurUnempTitle]
	# 	]
	
	# 1st Quarter Year Change
	if isYearChange:
		# Not needed since in table copied from MIMO spreadsheet
		# lSandRs.append([srLastEmpRateYears, srCurEmpRateYears])
		lSandRs.append([srLastYear, srCurYear])

	return curQtr, lastQtr, lSandRs

def menuFunction():
	lao.banner('MIMO Create New Qtr Files and Folders v01')
	msg = \
		'Select To Create New Qtr' \
		'\n\n  1) Market Insights' \
		'\n  2) Market Overviews' \
		'\n  3) Spreadsheets' \
		'\n  4) All' \
		'\n\n 00) Quit'
	print(msg)
	while 1:
		ui = td.uInput('\n ----- > ')
		if ui == '1':
			return ['Market Insights']
		elif ui == '2':
			return ['Market Overviews']
		elif ui == '3':
			return ['Spreadsheets']
		elif ui == '4':
			return ['Market Insights', 'Market Overviews', 'Spreadsheets']
		elif ui == '00':
			exit('\n Terminating program...')
		else:
			td.warningMsg('\n Invalid input...try again...')

def make_Market_Insigths_Folders_Files():

	# If current quarter directory does not exist make it
	if not os.path.exists('{0}{1}'.format(miFolder, curQtr)):
		os.mkdir('{0}{1}'.format(miFolder, curQtr))
	if not os.path.exists('{0}{1}/Eblast Covers'.format(miFolder, curQtr)):
		os.mkdir('{0}{1}/Eblast Covers'.format(miFolder, curQtr))
	if not os.path.exists('{0}{1}/PDFs'.format(miFolder, curQtr)):
		os.mkdir('{0}{1}/PDFs'.format(miFolder, curQtr))
	lINDD = glob('{0}{1}/*.indd'.format(miFolder, lastQtr))

	for row in lINDD:
		newINDDfile = row.replace(lastQtr, curQtr)
		print(' Copying: {0}'.format(newINDDfile))
		shutil.copy(row, newINDDfile)

def make_Market_Overviews_Folders_Files():
	# If current quarter directory does not exist make it
	if not os.path.exists('{0}{1}'.format(moFolder, curQtr)):
		os.mkdir('{0}{1}'.format(moFolder, curQtr))

	# Get file names of last quarter's PowerPionts.
	inFiles = glob('{0}{1}/*_v1.pptx'.format(moFolder, lastQtr))
	for fin in inFiles:
		print('\n {0}'.format(os.path.basename(fin)))

		# Rename PPTX file and directory
		# print(fin)
		skip_it = False
		fout = fin.replace(lastQtr, curQtr)
		fout = fout.replace('_v1', '')
		# Check if PowerPoint exists and if so let the user overwrite it
		if os.path.isfile(fout):
			td.warningMsg('\n PowerPoint exists...')
			while 1:
				ui = td.uInput('\n Overwrite the existing file [0/1/00] > ')
				if ui == '00':
					exit('\n Terminating program...')
				elif ui == '0':
					skip_it = True
					break
				elif ui == '1':
					skip_it = False
					break
				else:
					td.warningMsg('\n Invalid input...try again...')
		if skip_it is False:
			prs = Presentation(fin)
			# lSandRs = [['2nd', '1st'], ['March', 'June']
			for sr in lSandRs:
				# Loop through Master Slides
				for slide in prs.slide_layouts:
					loopSlideShapes(slide, sr[0], sr[1])
				# Loop through Slides
				for slide in prs.slides:
					loopSlideShapes(slide, sr[0], sr[1])
			print(' Saving...\n')
			prs.save(fout)

# Loop through slide's shapes and grouped shapes
def loopSlideShapes(slide, strSearch, strReplace):
	""""search and replace text in PowerPoint while preserving formatting"""
	#Useful Links ;)
	#https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
	#https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
	# Loop through groups of shapes
	group_shapes = [
			shp for shp in slide.shapes
			if shp.shape_type == MSO_SHAPE_TYPE.GROUP]

	for group_shape in group_shapes:
		for shape in group_shape.shapes:
			replaceTextInShape(shape, strSearch, strReplace)
	# Loop through shapes
	for shape in slide.shapes:
		replaceTextInShape(shape, strSearch, strReplace)

# Loop through shapes
def replaceTextInShape(shape, strSearch, strReplace):
	if shape.has_text_frame:
		# Determine if text to be replaced is in text frame
		if(shape.text.find(strSearch))!=-1:
			# print('\n FOUND {0}\n'.format(strSearch))
			text_frame = shape.text_frame
			# Get number of lines (carrage returns) in text frame
			lenPar = len(text_frame.paragraphs)
			# print(lenPar)
			i = 0
			# Loop through lines and replace text
			for i in range(0, lenPar):
				try:
					cur_text = text_frame.paragraphs[i].runs[0].text
					# print('\n Curent  text: {0}'.format(cur_text))
				except IndexError:
					# print(' FAILED: {0}'.format(strSearch))
					continue
				new_text = cur_text.replace(str(strSearch), str(strReplace))
				text_frame.paragraphs[i].runs[0].text = new_text

def make_Spreadsheets_Folders_Files():
	print('here spreadsheets')
	if not os.path.exists('{0}{1}'.format(sprshtFolder, curQtr)):
		os.mkdir('{0}{1}'.format(sprshtFolder, curQtr))
	if not os.path.exists('{0}{1}/Eblast Covers'.format(sprshtFolder, curQtr)):
		os.mkdir('{0}{1}/Eblast Covers'.format(sprshtFolder, curQtr))
	if not os.path.exists('{0}{1}/PDFs'.format(sprshtFolder, curQtr)):
		os.mkdir('{0}{1}/PDFs'.format(sprshtFolder, curQtr))
	lXLSX = glob('{0}{1}/*.xls*'.format(sprshtFolder, lastQtr))

	for row in lXLSX:
		if 'MIMO Data' in row:
			newXLSXfile = '{0}{1}/MIMO Data {1} v01.xlsm'.format(sprshtFolder, curQtr)
		else:
			newXLSXfile = row.replace(lastQtr, curQtr)
		print(' Copying: {0}'.format(newXLSXfile))
		shutil.copy(row, newXLSXfile)


curQtr, lastQtr, lSandRs = get_Qtrs_lSandRs()

# Variables
miOldFolder = 'F:/Research Department/MIMO/Market Insights/{0}'.format(lastQtr)
miNewFolder = 'F:/Research Department/MIMO/Market Insights/{0}'.format(curQtr)
miEblastCoversFolder = '{0}/Eblast Covers'.format(miNewFolder)
miPDFFolder = '{0}/PDFs'.format(miNewFolder)
miFolder = 'F:/Research Department/MIMO/Market Insights/'
moFolder = 'F:/Research Department/MIMO/PowerPoints/'
sprshtFolder = 'F:/Research Department/MIMO/Spreadsheets/'


while 1:
	lSelection = menuFunction()

	for sel in lSelection:
		while 1:
		
			lao.banner('Create New {0} Folders/Files for Quarter'.format(sel))
			
			print('\n Copying {0} {1} to {2}'.format(sel, lastQtr, curQtr))
			ui = td.uInput('\n Proceed [0/1/00] > ')
			if ui == '00':
				exit('\n Terminating program...')
			elif ui == '0': # Skipping
				sel = 'None'
				break
			elif ui == '1':
				break
			else:
				td.warningMsg('\n Invalid entry...try again...')
				lao.sleep(2)

		if sel == 'Market Insights':
			make_Market_Insigths_Folders_Files()
		elif sel == 'Market Overviews':
			make_Market_Overviews_Folders_Files()
		elif sel == 'Spreadsheets':
			make_Spreadsheets_Folders_Files()
	
	if len(lSelection) == 3:
		print('\n Copying completed...')
		exit('\n Terminating program...')

	ui = td.uInput('\n [Enter] Continue or [00] Quit > ')
	if ui == '00':
		exit('\n Terminating program...')
